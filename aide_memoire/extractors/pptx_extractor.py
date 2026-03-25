"""PowerPoint (.pptx) content extractor."""

from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from aide_memoire.extractors.base import BaseExtractor
from aide_memoire.models import ContentBlock, ContentType, ExtractedDocument


class PptxExtractor(BaseExtractor):
    def extract(self, file_path: Path, output_dir: Path | None = None) -> ExtractedDocument:
        prs = Presentation(str(file_path))
        doc = ExtractedDocument(source_path=file_path, title=file_path.stem)

        image_dir = None
        if output_dir:
            image_dir = output_dir / "images"
            image_dir.mkdir(parents=True, exist_ok=True)

        for slide_idx, slide in enumerate(prs.slides, 1):
            location = f"Slide {slide_idx}"

            for shape in slide.shapes:
                # Extract text from text frames
                if shape.has_text_frame:
                    text = self._extract_text_frame(shape.text_frame)
                    if text.strip():
                        doc.blocks.append(ContentBlock(
                            content_type=ContentType.TEXT,
                            text=text,
                            source_file=file_path.name,
                            source_location=location,
                        ))

                # Extract tables
                if shape.has_table:
                    table_text = self._extract_table(shape.table)
                    if table_text.strip():
                        doc.blocks.append(ContentBlock(
                            content_type=ContentType.TABLE,
                            text=table_text,
                            source_file=file_path.name,
                            source_location=location,
                        ))

                # Extract images
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and image_dir:
                    image_path = self._save_image(shape, image_dir, slide_idx)
                    if image_path:
                        doc.blocks.append(ContentBlock(
                            content_type=ContentType.IMAGE,
                            text=str(image_path),
                            source_file=file_path.name,
                            source_location=location,
                        ))

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_tf = slide.notes_slide.notes_text_frame
                if notes_tf and notes_tf.text.strip():
                    doc.blocks.append(ContentBlock(
                        content_type=ContentType.TEXT,
                        text=f"[Speaker Notes] {notes_tf.text}",
                        source_file=file_path.name,
                        source_location=location,
                    ))

        return doc

    def _extract_text_frame(self, text_frame) -> str:
        """Extract text from a text frame, preserving paragraph breaks."""
        paragraphs = []
        for para in text_frame.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
        return "\n".join(paragraphs)

    def _extract_table(self, table) -> str:
        """Extract table content as a structured text representation."""
        rows = []
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(" | ".join(cells))
        return "\n".join(rows)

    def _save_image(self, shape, image_dir: Path, slide_idx: int) -> Path | None:
        """Save an embedded image and return its path."""
        try:
            image = shape.image
            ext = image.content_type.split("/")[-1]
            if ext == "jpeg":
                ext = "jpg"
            filename = f"slide{slide_idx}_{shape.shape_id}.{ext}"
            image_path = image_dir / filename
            image_path.write_bytes(image.blob)
            return image_path
        except Exception:
            return None
